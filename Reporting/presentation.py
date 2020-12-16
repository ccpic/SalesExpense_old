from pptx import Presentation
from Reporting.reporting import *
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TOP = Inches(1.2)
LEFT_ONE_IMAGE = Inches(0.2)
WIDTH_ONE_IMAGE = Inches(13)


def add_sep_slide(title=None, layout_style=1, *args):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_style])

    title_placeholder = slide.placeholders[0]
    title_placeholder.text = title
    for i in range(2):
        font = title_placeholder.text_frame.paragraphs[i].runs[0].font
        font.size = Pt(32)

    print("Page%s" % str(int(prs.slides.index(slide)) + 1))


def add_img_slide(title=None, layout_style=0, *args):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_style])

    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title

    for item in args:
        img_path = item["img"]
        top = item["top"]
        left = item["left"]
        if "width" in item and "height" not in item:
            width = item["width"]
            slide.shapes.add_picture(img_path, left=left, top=top, width=width)
        elif "width" not in item and "height" in item:
            height = item["height"]
            slide.shapes.add_picture(img_path, left=left, top=top, height=height)
        elif "width" in item and "height" in item:
            width = item["width"]
            height = item["height"]
            slide.shapes.add_picture(img_path, left=left, top=top, width=width, height=height)

    print("Page%s" % str(int(prs.slides.index(slide)) + 1))


# 准备数据
df_pre = pd.read_excel("20201130095848.xlsx")  # 作为对比的上个时期的档案数据
df_post = pd.read_excel("20201215134151.xlsx")  # 作为对比的上个时期的档案数据
df_decile = pd.read_excel("decile.xlsx")  # 医院Decile数据文件，用于Decile相关分析的匹配

df_pre["月份"] = 202011
df_post["月份"] = 202012
df_total = pd.concat([df_pre, df_post])
df_total = pd.merge(df_total, df_decile.loc[:, ["医院编码", "IQVIA医院潜力", "IQVIA医院潜力分位"]], how="left", on="医院编码")
df_total = cleandata(df_total)

# 分南北中国
bu = "南中国"
df_pre = df_total[(df_total["南北中国"] == bu) & (df_total["月份"] == 202011)]
df_post = df_total[(df_total["南北中国"] == bu) & (df_total["月份"] == 202012)]
df_total = df_total[df_total["南北中国"] == bu]

pre = Clientfile(df_pre, name="南中国11月")
post = Clientfile(df_post, name="南中国12月")
total = Clientfile(df_total, name="南中国")

print("Data Ready")

# 创建ppt
prs = Presentation("template.pptx")
w = prs.slide_width
h = prs.slide_height


# Page3 标题分隔页
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.placeholders[0]
title.text = "现有客户档案\n数量和分布"
for i in range(2):
    font = title.text_frame.paragraphs[i].runs[0].font
    font.size = Pt(32)

print("Page3")

# Page4 客户档案基本情况
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.placeholders[0]
title.text = "%s客户档案基本情况" % post.name

data_list = [
    [df_post.shape[0], df_post.shape[0] - df_pre.shape[0], df_post.shape[0] / df_pre.shape[0] - 1,],
    [
        len(df_post["医院"].unique()),
        len(df_post["医院"].unique()) - len(df_pre["医院"].unique()),
        len(df_post["医院"].unique()) / len(df_pre["医院"].unique()) - 1,
    ],
    [
        len(df_post["地区经理"].unique()),
        len(df_post["地区经理"].unique()) - len(df_pre["地区经理"].unique()),
        len(df_post["地区经理"].unique()) / len(df_pre["地区经理"].unique()) - 1,
    ],
    [
        len(df_post["负责代表"].unique()),
        len(df_post["负责代表"].unique()) - len(df_pre["负责代表"].unique()),
        len(df_post["负责代表"].unique()) / len(df_pre["负责代表"].unique()) - 1,
    ],
]


shapes = slide.shapes
text_list = ["位医生", '家医院', '位地区经理', '位代表']
color_list = [RGBColor(0, 0, 128), RGBColor(220, 20, 60), RGBColor(255, 192, 0), RGBColor(29, 147, 60)]
for i in range(4):
    # 添加带数据的圆形
    oval_width = oval_height = Inches(2)
    shape = shapes.add_shape(
        MSO_SHAPE.OVAL, w / 8 * (2 * i + 1) - oval_width / 2, h / 2.3 - oval_height / 2, oval_width, oval_height
    )
    shape.fill.background()
    line = shape.line
    line.color.rgb = color_list[i]
    line.width = Pt(2.5)

    text_frame = shape.text_frame
    text_frame.word_wrap = False

    p = text_frame.paragraphs[0]
    p.text = "{:,.0f}".format(data_list[i][0])
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.size = Pt(28)

    p = text_frame.add_paragraph()
    p.text = "(%s, %s)" % ("{:+,.0f}".format(data_list[i][1]), "{:+.1%}".format(data_list[i][2]))
    p.alignment = PP_ALIGN.CENTER
    if data_list[i][1] > 0:
        p.font.color.rgb = RGBColor(0, 176, 80)
    elif data_list[i][1] < 0:
        p.font.color.rgb = RGBColor(255, 0, 0)
    else:
        p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.size = Pt(18)

    # 添加圆形下方标签
    textbox = shapes.add_textbox(w / 8 * (2 * i + 1) - oval_width / 2, h / 2.3 + oval_height / 2, oval_width, oval_height/2)
    text_frame =textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = text_list[i]
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = color_list[i]

print("Page4")

# Page 5 各医院/大区经理/代表客户档案覆盖数量分布
top = Inches(1.3)
height = Inches(5.1)
add_img_slide(
    "各医院/大区经理/代表客户档案覆盖数量分布",
    0,
    {
        "img": post.plot_hist_dist(
            pivoted=True,
            index="医院",
            bins=50,
            show_kde=True,
            show_tiles=False,
            show_metrics=True,
            xlim=[0, 200],
            width=6,
            height=8,
            pre=pre,
        ),
        "top": top,
        "left": Inches(0.2),
        "height": height,
    },
    {
        "img": post.plot_hist_dist(
            pivoted=True,
            index="地区经理",
            bins=10,
            show_kde=True,
            show_tiles=False,
            show_metrics=True,
            xlim=[0, 700],
            width=6,
            height=8,
            pre=pre,
        ),
        "top": top,
        "left": Inches(4.7),
        "height": height,
    },
    {
        "img": post.plot_hist_dist(
            pivoted=True,
            index="负责代表",
            bins=30,
            show_kde=True,
            show_tiles=False,
            show_metrics=True,
            xlim=[0, 200],
            width=6,
            height=8,
            pre=pre,
        ),
        "top": top,
        "left": Inches(8.8),
        "height": height,
    },
)

# Page6 客户档案基本分布情况
add_img_slide(
    "客户档案单位覆盖情况",
    0,
    {
        "img": total.plot_barline_dist(
            index="月份", columns="医院层级", values=None, perc=True, format_perc="{:.1%}", width=2, height=6
        ),
        "top": top,
        "left": Inches(0.2),
        "height": height,
    },
    {
        "img": total.plot_barline_dist(
            index="月份", columns="科室", values=None, perc=True, format_perc="{:.1%}", width=2, height=6
        ),
        "top": top,
        "left": Inches(4.7),
        "height": height,
    },
    {
        "img": total.plot_barline_dist(
            index="月份", columns="职称", values=None, perc=True, format_perc="{:.1%}", width=2, height=6
        ),
        "top": top,
        "left": Inches(8.8),
        "height": height,
    },
)

# Page7-12 档案数量相关指标汇总
index_list = ["医院层级", "IQVIA医院潜力分位", "科室", "职称", "区域", "大区"]
for idx in index_list:
    add_img_slide(
        "分%s档案数量相关指标汇总" % idx,
        0,
        {
            "img": post.plot_barh_kpi(index=idx, dimension="number", pre=pre),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )

# Page13-24 各大区多种维度级档案数分布
columns_list = ["医院层级", "IQVIA医院潜力分位", "科室", "职称"]
for col in columns_list:
    add_img_slide(
        "各大区分%s档案数" % col,
        0,
        {
            "img": post.plot_barline_dist(index="大区", columns=col, values=None, perc=False),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )
    add_img_slide(
        "各大区分%s档案数份额" % col,
        0,
        {
            "img": post.plot_barline_dist(index="大区", columns=col, values=None, perc=True),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )
    add_img_slide(
        "各大区分%s档案数变化" % col,
        0,
        {
            "img": post.plot_barline_dist(index="大区", columns=col, values=None, perc=False, pre=pre),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )

# Page25-27 各地区经理档案数量相关指标汇总
range_list = ((0, 17), (17, 34), (34, 51))
for i, r in enumerate(range_list):
    add_img_slide(
        "各地区经理档案数量相关指标汇总 - %s" % str(i + 1),
        0,
        {
            "img": post.plot_barh_kpi(
                index="地区经理", dimension="number", range=[r[0], r[1]], fontsize=12, mean_vline=True, pre=pre
            ),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )

# Page28 标题分隔页
add_sep_slide("现有客户档案\n潜力分析")

# Page29 档案数量十分位分布
add_img_slide(
    "客户档案潜力10分位划分",
    0,
    {
        "img": post.plot_hist_dist(pivoted=False, show_kde=True, show_tiles=True, bins=100, tiles=10),
        "top": TOP,
        "left": LEFT_ONE_IMAGE,
        "width": WIDTH_ONE_IMAGE
    },
)

# Page30 潜力饼图
height = Inches(6)
add_img_slide(
    "分潜力级别档案数量及病人数份额",
    0,
    {
        "img": total.plot_barline_dist(
            index="月份", columns="潜力级别", values=None, perc=True, format_perc="{:.1%}", width=6, height=6
        ),
        "top": TOP,
        "left": Inches(0.5),
        "height": height,
    },
    {
        "img": total.plot_barline_dist(
            index="月份", columns="潜力级别", values="月累计相关病人数", perc=True, format_perc="{:.1%}", width=6, height=6
        ),
        "top": TOP,
        "left": Inches(7.5),
        "height": height,
    },
)

# Page31-36 档案潜力相关指标汇总
index_list = ["医院层级", "IQVIA医院潜力分位", "科室", "职称", "区域", "大区"]
for idx in index_list:
    add_img_slide(
        "分%s档案潜力相关指标汇总" % idx,
        0,
        {
            "img": post.plot_barh_kpi(index=idx, dimension="potential"),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )

# Page37-39 各地区经理档案数量相关指标汇总
for i, r in enumerate(range_list):
    add_img_slide(
        "各地区经理潜力数量相关指标汇总 - %s" % str(i + 1),
        0,
        {
            "img": post.plot_barh_kpi(
                index="地区经理", dimension="potential", range=[r[0], r[1]], fontsize=12, mean_vline=True
            ),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )

# Page 40 标题间隔页
add_sep_slide("数量 vs. 潜力\n交叉分析")

# Page 41-44 数量潜力气泡图
index_list = ["大区", "地区经理"]
for idx in index_list:
    add_img_slide(
        "%s客户档案数 versus 平均潜力" % idx,
        0,
        {
            "img": post.plot_bubble_number_potential(idx, z_scale=0.02, labelLimit=100),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )
    add_img_slide(
        "%s客户档案数（平均每代表） versus 平均潜力" % idx,
        0,
        {
            "img": post.plot_bubble_number_potential(idx, z_scale=0.07, labelLimit=100, dimension="rsp"),
            "top": TOP,
            "left": LEFT_ONE_IMAGE,
            "width": WIDTH_ONE_IMAGE
        },
    )


# 保存ppt文件
prs.save("presentation.pptx")
print("PPT has been saved")
